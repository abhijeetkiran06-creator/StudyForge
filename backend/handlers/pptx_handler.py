import os
from pptx import Presentation


def extract_pptx(file_path):

    presentation = Presentation(file_path)

    text = []

    for slide in presentation.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):
                text.append(shape.text)

    return {
        "filename": os.path.basename(file_path),
        "extension": ".pptx",
        "text": "\n".join(text)
    }